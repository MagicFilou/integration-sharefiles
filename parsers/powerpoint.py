from pptx import Presentation
from parsers import Parser


class PowerPointParser(Parser):

    ACCEPTED_TYPES = ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt",".pptx"]

    def parse(self, item):
        super().parse(item)

        #item is a Path or file
        pres = Presentation(item)

        slide_no = 0
        for slide in pres.slides:
            shape_no = 0
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # We assume this is a filepath we can read from
                    value = shape.text.strip()
                    if value:
                        yield (
                            shape.text,
                            f"{slide_no}_{shape_no}"
                            # {
                            #     "type": "powerpoint",
                            #     "position":{
                            #         "slide": slide_no
                            #     }
                            # }
                        )
                shape_no +=1
            slide_no += 1


if __name__ == "__main__":
    import sys
    parser = PowerPointParser(None) # Fake the source for now.
    for data in parser.parse(sys.argv[1]):
        text, meta = data
        print(data)

